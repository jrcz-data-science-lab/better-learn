from pdfminer.high_level import extract_text
from pdfminer.layout import LAParams
from pptx import Presentation

def extract_pdf_text(path):
    laparams = LAParams(char_margin=1.5, line_margin=0.8, word_margin=0.1, boxes_flow=0.5)
    text = extract_text(path, laparams=laparams)
    
    if text and not text.endswith('\n'):
        text += '\n'
    
    lines = [line for line in text.splitlines() if line.strip()]
    
    if lines and not lines[-1].endswith(('.', '?', '!', '…')):
        lines[-1] += '.'
    
    return "\n".join(lines)

def extract_pptx_text(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        shapes = [s for s in slide.shapes if s.has_text_frame]
        shapes.sort(key=lambda s: (s.top, s.left))
        for shape in shapes:
            line = shape.text.strip()
            if line:
                texts.append(line)
    return "\n".join(texts)

if __name__ == "__main__":
    pdf_text = extract_pdf_text("./sample_files/sample.pdf")
    print("PDF Content:\n", pdf_text, "\n")

    pptx_text = extract_pptx_text("./sample_files/sample.pptx")
    print("PPTX Content:\n", pptx_text, "\n")